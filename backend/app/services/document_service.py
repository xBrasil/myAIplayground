"""Extract plain text from PDF, DOCX, XLSX, and PPTX files."""

import logging
from io import BytesIO
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text(file_path: str, max_chars: int = 256_000) -> str:
    """Return extracted text from a document file, truncated to *max_chars*."""
    suffix = Path(file_path).suffix.lower()
    try:
        if suffix == ".pdf":
            return _extract_pdf(file_path, max_chars)
        if suffix == ".docx":
            return _extract_docx(file_path, max_chars)
        if suffix == ".xlsx":
            return _extract_xlsx(file_path, max_chars)
        if suffix == ".pptx":
            return _extract_pptx(file_path, max_chars)
    except Exception as exc:
        logger.warning("Falha ao extrair texto de %s: %s", file_path, exc)
        return f"(Erro ao extrair texto do documento: {exc})"
    return "(Formato de documento não reconhecido)"


def _extract_pdf(path: str, max_chars: int) -> str:
    import pymupdf  # lazy import

    parts: list[str] = []
    total = 0
    with pymupdf.open(path) as doc:
        for page in doc:
            text = page.get_text()
            parts.append(text)
            total += len(text)
            if total >= max_chars:
                break
    return "".join(parts)[:max_chars]


def _extract_docx(path: str, max_chars: int) -> str:
    from docx import Document  # lazy import

    doc = Document(path)
    parts: list[str] = []
    total = 0
    for para in doc.paragraphs:
        parts.append(para.text)
        total += len(para.text) + 1
        if total >= max_chars:
            break
    return "\n".join(parts)[:max_chars]


def _extract_xlsx(path: str, max_chars: int) -> str:
    from openpyxl import load_workbook  # lazy import

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    total = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Planilha: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            parts.append(row_text)
            total += len(row_text) + 1
            if total >= max_chars:
                break
        if total >= max_chars:
            break
    wb.close()
    return "\n".join(parts)[:max_chars]


def _extract_pptx(path: str, max_chars: int) -> str:
    from pptx import Presentation  # lazy import

    prs = Presentation(path)
    parts: list[str] = []
    total = 0
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text
                    parts.append(text)
                    total += len(text) + 1
                    if total >= max_chars:
                        break
            if total >= max_chars:
                break
        if total >= max_chars:
            break
    return "\n".join(parts)[:max_chars]
